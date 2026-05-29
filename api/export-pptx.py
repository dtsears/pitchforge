"""
Vercel Python serverless function.
POST /api/export-pptx
Body: { slides: Slide[], prospectName: string }
Returns: .pptx binary
"""

from http.server import BaseHTTPRequestHandler
import json
import io
import os
import sys

# ── Template slide indices (0-based) ────────────────────────────────────────
# Mapped from the Bluehost 2026 template analysis
SLIDE_MAP = {
    "COVER":         1,   # Slide 2  — dark left panel + headline right
    "REP":           5,   # Slide 6  — panelist layout with photo
    "UNDERSTANDING": 23,  # Slide 24 — title + body + right photo panel
    "WHY_NOW":       20,  # Slide 21 — left text + right photo panel
    "SOLUTION":      41,  # Slide 42 — title + 4-row icon list
    "PROOF":         26,  # Slide 27 — headline + stat cards
    "ROI":           19,  # Slide 20 — left body + right callout box
    "NEXT_STEPS":    17,  # Slide 18 — two-column text
}

DECK_ORDER = ["COVER", "REP", "UNDERSTANDING", "WHY_NOW", "SOLUTION", "PROOF", "ROI", "NEXT_STEPS"]


def find_template():
    candidates = [
        "/var/task/public/bluehost-template.pptx",
        os.path.join(os.path.dirname(__file__), "..", "public", "bluehost-template.pptx"),
        os.path.join(os.getcwd(), "public", "bluehost-template.pptx"),
    ]
    for p in candidates:
        if os.path.exists(p):
            return p
    raise FileNotFoundError(f"Template not found. Tried: {candidates}")


# ── Text replacement ─────────────────────────────────────────────────────────

def set_text(tf, text: str):
    """
    Replace all text in a TextFrame with `text`, preserving the
    font/size/color of the very first run on the slide so the
    Bluehost branding is kept intact.
    """
    if not tf.paragraphs:
        return

    first_para = tf.paragraphs[0]
    first_run = first_para.runs[0] if first_para.runs else None

    # Clear every paragraph
    from pptx.oxml.ns import qn
    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    # Re-add one paragraph per line
    lines = text.split("\n") if text else [""]
    for line_idx, line in enumerate(lines):
        import copy
        if first_para._p is not None:
            new_p = copy.deepcopy(first_para._p)
        else:
            from lxml import etree
            new_p = etree.SubElement(txBody, qn("a:p"))

        # Remove existing runs in cloned para
        for r in new_p.findall(qn("a:r")):
            new_p.remove(r)

        # Build a run
        if first_run is not None:
            new_r = copy.deepcopy(first_run._r)
            t_el = new_r.find(qn("a:t"))
            if t_el is not None:
                t_el.text = line
            else:
                from lxml import etree
                t_el = etree.SubElement(new_r, qn("a:t"))
                t_el.text = line
            new_p.append(new_r)
        else:
            from lxml import etree
            r_el = etree.SubElement(new_p, qn("a:r"))
            t_el = etree.SubElement(r_el, qn("a:t"))
            t_el.text = line

        txBody.append(new_p)


def shapes_with_text(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def largest_shape(slide):
    """Return the shape whose first run has the largest font size."""
    best, best_size = None, 0
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        for para in s.text_frame.paragraphs:
            for run in para.runs:
                sz = run.font.size or 0
                if sz > best_size:
                    best_size = sz
                    best = s
    return best


# ── Slide populators ─────────────────────────────────────────────────────────

def pop_cover(slide, content, prospect_name):
    heading = largest_shape(slide)
    if heading:
        set_text(heading.text_frame, content.get("headline", ""))
    for s in shapes_with_text(slide):
        if s == heading:
            continue
        t = s.text_frame.text.lower()
        if any(k in t for k in ["present", "name &", "heading", "copy"]):
            set_text(s.text_frame, f"Prepared for {prospect_name}")
            break


def pop_rep(slide, content, prospect_name):
    heading = largest_shape(slide)
    if heading:
        set_text(heading.text_frame, content.get("name", "Your Rep"))

    specialties = ", ".join(content.get("specialties", []))
    body_text = "\n".join([
        content.get("title", ""),
        f"{content.get('tenureYears', 0)} years at Bluehost",
        content.get("bio", ""),
        specialties,
    ])
    for s in shapes_with_text(slide):
        if s == heading:
            continue
        if len(s.text_frame.paragraphs) >= 2 or len(s.text_frame.text) > 30:
            set_text(s.text_frame, body_text)
            break


def pop_understanding(slide, content, prospect_name):
    heading = largest_shape(slide)
    if heading:
        set_text(heading.text_frame, content.get("title", "What We Heard"))

    pain_points = content.get("painPoints", [])
    lines = []
    for pp in pain_points:
        lines.append(f"• {pp.get('headline', '')}")
        lines.append(f"  {pp.get('description', '')}")
    body_text = "\n".join(lines)

    for s in shapes_with_text(slide):
        if s == heading:
            continue
        set_text(s.text_frame, body_text)
        break


def pop_why_now(slide, content, prospect_name):
    heading = largest_shape(slide)
    if heading:
        set_text(heading.text_frame, content.get("title", "Why Now"))

    body_lines = [
        content.get("triggerEvent", ""),
        "",
        content.get("description", ""),
        "",
    ] + [f"→ {pt}" for pt in content.get("urgencyPoints", [])]

    for s in shapes_with_text(slide):
        if s == heading:
            continue
        set_text(s.text_frame, "\n".join(body_lines))
        break


def pop_solution(slide, content, prospect_name):
    heading = largest_shape(slide)
    if heading:
        set_text(heading.text_frame, content.get("title", "How Bluehost Helps"))

    products = content.get("products", [])
    text_shapes = [s for s in shapes_with_text(slide) if s != heading]

    # Try to set each product into its own shape
    for i, product in enumerate(products[:4]):
        if i < len(text_shapes):
            tf = text_shapes[i].text_frame
            lines = [product.get("name", ""), product.get("benefit", "")]
            set_text(tf, "\n".join(lines))
        # If we run out of shapes, combine into one
    if not text_shapes and products:
        pass  # Nothing to populate in this case
    elif len(products) > len(text_shapes) and text_shapes:
        combined = "\n\n".join(
            f"{p.get('name', '')}\n{p.get('benefit', '')}" for p in products
        )
        set_text(text_shapes[0].text_frame, combined)


def pop_proof(slide, content, prospect_name):
    heading = largest_shape(slide)
    if heading:
        set_text(heading.text_frame, content.get("title", "Results That Speak"))

    metric = content.get("headlineMetric", "")
    narrative = content.get("narrative", "")
    customer = content.get("customerName", "")
    source_url = content.get("sourceUrl", "")
    products_used = " · ".join(content.get("productsUsed", []))

    body_text = "\n".join(filter(None, [
        narrative,
        "",
        f"— {customer}" if customer else "",
        products_used,
        source_url,
    ]))

    shapes = [s for s in shapes_with_text(slide) if s != heading]
    # Largest remaining shape gets the metric; next gets narrative
    shapes_sorted = sorted(shapes, key=lambda s: s.text_frame.paragraphs[0].runs[0].font.size or 0, reverse=True)
    for i, s in enumerate(shapes_sorted[:3]):
        if i == 0:
            set_text(s.text_frame, metric)
        else:
            set_text(s.text_frame, body_text)
            break


def pop_roi(slide, content, prospect_name):
    heading = largest_shape(slide)
    if heading:
        set_text(heading.text_frame, content.get("title", "The Business Case"))

    metrics = content.get("metrics", [])
    summary = content.get("summary", "")

    shapes = [s for s in shapes_with_text(slide) if s != heading]
    for i, s in enumerate(shapes[:4]):
        if i == 0:
            set_text(s.text_frame, summary)
        elif i - 1 < len(metrics):
            m = metrics[i - 1]
            set_text(s.text_frame, f"{m.get('value', '')}\n{m.get('label', '')}\n{m.get('description', '')}")


def pop_next_steps(slide, content, prospect_name):
    heading = largest_shape(slide)
    if heading:
        set_text(heading.text_frame, content.get("title", "Next Steps"))

    steps = content.get("steps", [])
    cta = content.get("cta", "")
    step_lines = [
        f"{i + 1}. {s.get('action', '')}  |  {s.get('owner', '')}  |  {s.get('timeline', '')}"
        for i, s in enumerate(steps)
    ]

    shapes = [s for s in shapes_with_text(slide) if s != heading]
    if shapes:
        set_text(shapes[0].text_frame, "\n".join(step_lines))
    if len(shapes) > 1:
        set_text(shapes[1].text_frame, cta)


POPULATORS = {
    "COVER":         pop_cover,
    "REP":           pop_rep,
    "UNDERSTANDING": pop_understanding,
    "WHY_NOW":       pop_why_now,
    "SOLUTION":      pop_solution,
    "PROOF":         pop_proof,
    "ROI":           pop_roi,
    "NEXT_STEPS":    pop_next_steps,
}


# ── Slide deletion ───────────────────────────────────────────────────────────

def remove_slide(prs, index: int):
    """Remove a slide from the presentation by 0-based index."""
    slide_id_list = prs.slides._sldIdLst
    if index >= len(slide_id_list):
        return
    slide_el = slide_id_list[index]
    r_id = slide_el.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    slide_id_list.remove(slide_el)
    if r_id:
        try:
            prs.part.drop_rel(r_id)  # python-pptx 1.0+ API
        except Exception:
            pass  # relationship may already be gone


def keep_slides(prs, indices_to_keep: list):
    """Remove all slides NOT in indices_to_keep. Delete from the end."""
    total = len(prs.slides._sldIdLst)
    to_remove = sorted(
        [i for i in range(total) if i not in indices_to_keep],
        reverse=True,
    )
    for i in to_remove:
        remove_slide(prs, i)


# ── HTTP handler ─────────────────────────────────────────────────────────────

def make_filename(prospect_name: str, website_url: str) -> str:
    """Bluehost_CompanyName_domain.com.pptx"""
    try:
        from urllib.parse import urlparse
        domain = urlparse(website_url).netloc or website_url
        domain = domain.replace("www.", "").split("/")[0]
    except Exception:
        domain = ""
    safe_name = "".join(c if c.isalnum() or c in "-_" else "_" for c in prospect_name)
    safe_domain = "".join(c if c.isalnum() or c in "-_." else "_" for c in domain)
    parts = ["Bluehost", safe_name]
    if safe_domain:
        parts.append(safe_domain)
    return "_".join(parts) + ".pptx"


class handler(BaseHTTPRequestHandler):

    def do_OPTIONS(self):
        self.send_response(200)
        self._cors()
        self.end_headers()

    def do_POST(self):
        try:
            length = int(self.headers.get("Content-Length", 0))
            body = json.loads(self.rfile.read(length))

            slides_data = body.get("slides", [])
            prospect_name = body.get("prospectName", "Prospect")
            website_url = body.get("websiteUrl", "")

            prs = Presentation(find_template())

            # Populate slides BEFORE deleting, using original indices
            for slide_item in slides_data:
                slide_type = slide_item.get("type", "")
                if slide_type not in SLIDE_MAP:
                    continue
                idx = SLIDE_MAP[slide_type]
                if idx < len(prs.slides):
                    populator = POPULATORS.get(slide_type)
                    if populator:
                        try:
                            populator(prs.slides[idx], slide_item.get("content", {}), prospect_name)
                        except Exception:
                            pass  # Never fail the whole export for one bad slide

            # Keep only the slides we need
            deck_types = [s["type"] for s in slides_data if s["type"] in SLIDE_MAP]
            keep_indices = [SLIDE_MAP[t] for t in deck_types]
            keep_slides(prs, keep_indices)

            buf = io.BytesIO()
            prs.save(buf)
            pptx_bytes = buf.getvalue()

            filename = make_filename(prospect_name, website_url)

            self.send_response(200)
            self._cors()
            self.send_header(
                "Content-Type",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            self.send_header("Content-Disposition", f'attachment; filename="{filename}"')
            self.send_header("Content-Length", str(len(pptx_bytes)))
            self.end_headers()
            self.wfile.write(pptx_bytes)

        except Exception as e:
            msg = str(e).encode()
            self.send_response(500)
            self._cors()
            self.send_header("Content-Type", "text/plain")
            self.send_header("Content-Length", str(len(msg)))
            self.end_headers()
            self.wfile.write(msg)

    def _cors(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")

    def log_message(self, *_):
        pass
